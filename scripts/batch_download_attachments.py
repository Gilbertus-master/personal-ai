"""
Batch download and import attachments for all corporate emails.

Queries the DB for company_email documents, checks Graph API for attachments,
downloads and extracts text (PDF, DOCX, PPTX, plain text), and imports into DB.

Uses retry logic with backoff for SSL/network resilience.
Tracks progress in a state file so runs can be resumed.

Usage:
    .venv/bin/python scripts/batch_download_attachments.py
    .venv/bin/python scripts/batch_download_attachments.py --limit 20
    .venv/bin/python scripts/batch_download_attachments.py --limit 20 --reset
"""
from __future__ import annotations

import argparse
import base64
import io
import json
import logging
import re
import sys
import time
from datetime import datetime
from html import unescape
from pathlib import Path
from typing import Any

import requests
from dotenv import load_dotenv

# Ensure project root is on sys.path so app.* imports work
PROJECT_ROOT = Path(__file__).resolve().parents[1]
sys.path.insert(0, str(PROJECT_ROOT))

load_dotenv(PROJECT_ROOT / ".env")

from app.ingestion.graph_api.auth import get_access_token
from app.ingestion.common.db import (
    document_exists_by_raw_path,
    insert_chunk,
    insert_document,
    insert_source,
)
from app.db.postgres import get_pg_connection

# ---------------------------------------------------------------------------
# Config
# ---------------------------------------------------------------------------

GRAPH_BASE = "https://graph.microsoft.com/v1.0"
STATE_FILE = PROJECT_ROOT / ".batch_attachments_state.json"

MAX_RETRIES = 3
BACKOFF_SECS = 10
REQUEST_TIMEOUT_LIST = 30
REQUEST_TIMEOUT_DOWNLOAD = 120

CHUNK_TARGET_CHARS = 2500
CHUNK_OVERLAP_CHARS = 250
ATTACHMENT_MAX_CHARS = 50_000

_NETWORK_ERRORS: tuple = (
    requests.RequestException,
    ConnectionError,
    OSError,
)

_SKIP_EXTENSIONS = frozenset((
    ".png", ".jpg", ".jpeg", ".gif", ".bmp", ".svg", ".ico", ".tif", ".tiff",
    ".webp", ".zip", ".rar", ".7z", ".gz", ".tar", ".exe", ".dll", ".msi",
    ".mp3", ".mp4", ".wav", ".avi", ".mov", ".wmv", ".flv", ".mkv",
    ".bin", ".dat", ".iso", ".img",
))

_TEXT_EXTENSIONS = frozenset((
    ".txt", ".csv", ".md", ".html", ".htm", ".xml", ".json", ".log", ".yaml",
    ".yml", ".ini", ".cfg", ".conf", ".tsv",
))

# ---------------------------------------------------------------------------
# Logging
# ---------------------------------------------------------------------------

logging.basicConfig(
    level=logging.INFO,
    format="%(asctime)s [%(levelname)s] %(message)s",
    datefmt="%H:%M:%S",
)
log = logging.getLogger("batch_attachments")

# ---------------------------------------------------------------------------
# State tracking (resume support)
# ---------------------------------------------------------------------------


def _load_state() -> dict[str, Any]:
    if STATE_FILE.exists():
        try:
            return json.loads(STATE_FILE.read_text())
        except (json.JSONDecodeError, OSError):
            pass
    return {"processed_msg_ids": [], "stats": {"checked": 0, "with_attachments": 0, "imported": 0, "skipped": 0, "errors": 0}}


def _save_state(state: dict[str, Any]) -> None:
    STATE_FILE.write_text(json.dumps(state, indent=2))


# ---------------------------------------------------------------------------
# Text extraction
# ---------------------------------------------------------------------------


def _extract_text_from_bytes(content: bytes, name: str) -> str | None:
    """Extract text from attachment bytes based on file extension."""
    lower = name.lower()

    # Skip known binary / image types
    for ext in _SKIP_EXTENSIONS:
        if lower.endswith(ext):
            return None

    if lower.endswith(".pdf"):
        try:
            from pypdf import PdfReader
            reader = PdfReader(io.BytesIO(content))
            text = "\n\n".join(p.extract_text() or "" for p in reader.pages)
            return text if text.strip() else None
        except Exception as e:
            log.warning("PDF extraction failed for %s: %s", name, e)
            return None

    if lower.endswith(".docx"):
        try:
            from docx import Document
            doc = Document(io.BytesIO(content))
            text = "\n".join(p.text for p in doc.paragraphs if p.text.strip())
            return text if text.strip() else None
        except Exception as e:
            log.warning("DOCX extraction failed for %s: %s", name, e)
            return None

    if lower.endswith(".pptx"):
        try:
            from pptx import Presentation
            prs = Presentation(io.BytesIO(content))
            parts: list[str] = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        for paragraph in shape.text_frame.paragraphs:
                            text = paragraph.text.strip()
                            if text:
                                parts.append(text)
            return "\n".join(parts) if parts else None
        except ImportError:
            log.warning("python-pptx not installed — skipping %s", name)
            return None
        except Exception as e:
            log.warning("PPTX extraction failed for %s: %s", name, e)
            return None

    if lower.endswith(".xlsx") or lower.endswith(".xls"):
        # Excel files need openpyxl — skip gracefully if unavailable
        try:
            from openpyxl import load_workbook
            wb = load_workbook(io.BytesIO(content), read_only=True, data_only=True)
            parts: list[str] = []
            for ws in wb.worksheets:
                for row in ws.iter_rows(values_only=True):
                    cells = [str(c) for c in row if c is not None]
                    if cells:
                        parts.append(" | ".join(cells))
            wb.close()
            return "\n".join(parts) if parts else None
        except ImportError:
            log.warning("openpyxl not installed — skipping %s", name)
            return None
        except Exception as e:
            log.warning("Excel extraction failed for %s: %s", name, e)
            return None

    # Plain text types
    for ext in _TEXT_EXTENSIONS:
        if lower.endswith(ext):
            try:
                text = content.decode("utf-8", errors="ignore")
                if lower.endswith((".html", ".htm")):
                    text = re.sub(r"(?is)<.*?>", " ", text)
                    text = unescape(text)
                return text if text.strip() else None
            except Exception as e:
                log.warning("Text extraction failed for %s: %s", name, e)
                return None

    # Generic fallback — try as text
    try:
        text = content.decode("utf-8", errors="ignore")
        if len(text.strip()) > 100:
            return text
    except Exception:
        pass

    return None


# ---------------------------------------------------------------------------
# Chunking
# ---------------------------------------------------------------------------


def _chunk_text(text: str) -> list[str]:
    text = (text or "").strip()
    if not text:
        return [""]

    chunks: list[str] = []
    start = 0
    length = len(text)
    while start < length:
        end = min(start + CHUNK_TARGET_CHARS, length)
        chunk = text[start:end].strip()
        if chunk:
            chunks.append(chunk)
        if end >= length:
            break
        start = max(end - CHUNK_OVERLAP_CHARS, start + 1)
    return chunks


# ---------------------------------------------------------------------------
# Graph API helpers with retry
# ---------------------------------------------------------------------------


def _graph_get_with_retry(
    url: str,
    token: str,
    params: dict | None = None,
    timeout: int = 30,
    label: str = "",
) -> dict[str, Any] | None:
    """GET request to Graph API with retry logic. Returns None on total failure."""
    headers = {"Authorization": f"Bearer {token}", "Accept": "application/json"}
    for attempt in range(MAX_RETRIES):
        try:
            resp = requests.get(url, headers=headers, params=params, timeout=timeout)
            resp.raise_for_status()
            return resp.json()
        except _NETWORK_ERRORS as e:
            log.warning(
                "%s attempt %d/%d failed: %s",
                label or url[:80], attempt + 1, MAX_RETRIES, e,
            )
            if attempt < MAX_RETRIES - 1:
                time.sleep(BACKOFF_SECS * (attempt + 1))
        except requests.HTTPError as e:
            # 404 = message deleted, 401 = token expired — don't retry
            status = getattr(e.response, "status_code", 0)
            if status in (401, 403):
                log.error("Auth error (%d) — token may be expired. Aborting.", status)
                raise
            if status == 404:
                log.warning("Resource not found (404) for %s — skipping", label)
                return None
            log.warning(
                "%s attempt %d/%d HTTP %d: %s",
                label or url[:80], attempt + 1, MAX_RETRIES, status, e,
            )
            if attempt < MAX_RETRIES - 1:
                time.sleep(BACKOFF_SECS * (attempt + 1))
    log.error("All %d retries exhausted for %s", MAX_RETRIES, label or url[:80])
    return None


# ---------------------------------------------------------------------------
# DB queries
# ---------------------------------------------------------------------------


def _get_email_documents() -> list[dict[str, Any]]:
    """
    Get all company_email documents from DB.
    Returns list of dicts with id, raw_path, title, author, created_at.
    """
    sql = """
        SELECT d.id, d.raw_path, d.title, d.author, d.created_at
        FROM documents d
        JOIN sources s ON d.source_id = s.id
        WHERE s.source_type = 'company_email'
        ORDER BY d.created_at ASC NULLS LAST
    """
    with get_pg_connection() as conn:
        with conn.cursor() as cur:
            cur.execute(sql)
            rows = cur.fetchall()

    results = []
    for row in rows:
        results.append({
            "id": row[0],
            "raw_path": row[1] or "",
            "title": row[2] or "",
            "author": row[3] or "",
            "created_at": row[4],
        })
    return results


def _get_existing_attachment_msg_ids() -> set[str]:
    """
    Get the set of Graph message IDs that already have attachment documents.
    Attachment raw_paths look like: graph://attachment/{msg_id}/{filename}
    """
    sql = """
        SELECT DISTINCT d.raw_path
        FROM documents d
        JOIN sources s ON d.source_id = s.id
        WHERE s.source_type = 'company_email_attachment'
    """
    with get_pg_connection() as conn:
        with conn.cursor() as cur:
            cur.execute(sql)
            rows = cur.fetchall()

    msg_ids: set[str] = set()
    for row in rows:
        raw_path = row[0] or ""
        # graph://attachment/{msg_id}/{filename}
        if raw_path.startswith("graph://attachment/"):
            parts = raw_path.split("/", 4)  # ['graph:', '', 'attachment', msg_id, filename]
            if len(parts) >= 4:
                msg_ids.add(parts[3])
    return msg_ids


def _extract_graph_msg_id(raw_path: str) -> str | None:
    """
    Extract the Graph message ID from a document raw_path.
    Patterns seen:
      graph://corporate_email/inbox/{msg_id}
      graph://priority_email/{msg_id}
    """
    if not raw_path or not raw_path.startswith("graph://"):
        return None
    # Take the last path segment as the msg_id (the long base64-ish ID)
    parts = raw_path.rstrip("/").split("/")
    if len(parts) >= 4:
        candidate = parts[-1]
        # Graph message IDs are long (>50 chars typically)
        if len(candidate) > 40:
            return candidate
    return None


# ---------------------------------------------------------------------------
# Main processing
# ---------------------------------------------------------------------------


def _process_email(
    doc: dict[str, Any],
    msg_id: str,
    token: str,
    att_source_id: int,
) -> int:
    """
    Check a single email for attachments and import any found.
    Returns number of attachment documents imported.
    """
    subject = doc["title"]
    author = doc["author"] or ""

    # 1. List attachments (metadata only — lightweight)
    url = f"{GRAPH_BASE}/me/messages/{msg_id}/attachments"
    data = _graph_get_with_retry(
        url, token,
        params={"$select": "id,name,size,contentType"},
        timeout=REQUEST_TIMEOUT_LIST,
        label=f"list-att {subject[:40]}",
    )
    if data is None:
        return 0

    att_list = data.get("value", [])
    if not att_list:
        return 0

    log.info(
        "  Email has %d attachment(s): %s",
        len(att_list), subject[:60],
    )

    received = doc.get("created_at")
    received_str = ""
    if received:
        if isinstance(received, str):
            received_str = received[:16]
        else:
            received_str = received.isoformat()[:16]

    imported = 0

    for att_meta in att_list:
        att_id = att_meta.get("id", "")
        att_name = att_meta.get("name", "unknown")
        content_type = att_meta.get("contentType", "")

        # Skip images and known binary content types
        if content_type and content_type.startswith(("image/", "video/", "audio/")):
            log.info("    Skipping %s (content-type: %s)", att_name, content_type)
            continue

        lower_name = att_name.lower()
        if any(lower_name.endswith(ext) for ext in _SKIP_EXTENSIONS):
            log.info("    Skipping %s (binary/image extension)", att_name)
            continue

        raw_path = f"graph://attachment/{msg_id}/{att_name}"
        if document_exists_by_raw_path(raw_path):
            log.info("    Already imported: %s", att_name)
            continue

        # 2. Download individual attachment content
        att_url = f"{GRAPH_BASE}/me/messages/{msg_id}/attachments/{att_id}"
        att_data = _graph_get_with_retry(
            att_url, token,
            timeout=REQUEST_TIMEOUT_DOWNLOAD,
            label=f"download {att_name}",
        )
        if att_data is None:
            continue

        content_b64 = att_data.get("contentBytes", "")
        if not content_b64:
            log.info("    No contentBytes for %s (may be reference attachment)", att_name)
            continue

        try:
            content_bytes = base64.b64decode(content_b64)
        except Exception as e:
            log.warning("    base64 decode failed for %s: %s", att_name, e)
            continue

        # 3. Extract text
        text = _extract_text_from_bytes(content_bytes, att_name)
        if not text or not text.strip():
            log.info("    No extractable text from %s", att_name)
            continue

        # 4. Clean and truncate
        clean = text.replace("\x00", "").strip()
        if len(clean) > ATTACHMENT_MAX_CHARS:
            clean = clean[:ATTACHMENT_MAX_CHARS] + "\n[...truncated]"

        full_text = (
            f"Attachment: {att_name}\n"
            f"From email: {subject}\n"
            f"Sender: {author}\n"
            f"Date: {received_str}\n\n"
            f"{clean}"
        )

        # 5. Import as document + chunks
        received_dt = received
        if isinstance(received_dt, str):
            try:
                received_dt = datetime.fromisoformat(received_dt.replace("Z", "+00:00"))
            except ValueError:
                received_dt = None

        document_id = insert_document(
            conn=None,
            source_id=att_source_id,
            title=f"[ATT] {att_name} — {subject[:40]}",
            created_at=received_dt,
            author=author,
            participants=[author] if author else [],
            raw_path=raw_path,
        )

        chunks = _chunk_text(full_text)
        for chunk_index, chunk in enumerate(chunks):
            insert_chunk(
                conn=None,
                document_id=document_id,
                chunk_index=chunk_index,
                text=chunk,
                timestamp_start=received_dt,
                timestamp_end=received_dt,
                embedding_id=None,
            )

        imported += 1
        log.info(
            "    Imported: %s (%d chunks, %d chars)",
            att_name, len(chunks), len(clean),
        )

    return imported


def main() -> None:
    parser = argparse.ArgumentParser(
        description="Batch download and import attachments for all corporate emails.",
    )
    parser.add_argument(
        "--limit", type=int, default=None,
        help="Max number of emails to process (for incremental runs)",
    )
    parser.add_argument(
        "--reset", action="store_true",
        help="Reset state file and reprocess all emails",
    )
    parser.add_argument(
        "--dry-run", action="store_true",
        help="Only check for attachments, don't import",
    )
    args = parser.parse_args()

    # Load or reset state
    if args.reset and STATE_FILE.exists():
        STATE_FILE.unlink()
        log.info("State file reset.")

    state = _load_state()
    already_processed: set[str] = set(state.get("processed_msg_ids", []))
    stats = state.get("stats", {"checked": 0, "with_attachments": 0, "imported": 0, "skipped": 0, "errors": 0})

    log.info("=" * 60)
    log.info("BATCH ATTACHMENT DOWNLOADER")
    log.info("=" * 60)

    # 1. Get all corporate email documents from DB
    email_docs = _get_email_documents()
    log.info("Found %d company_email documents in DB", len(email_docs))

    # 2. Get set of msg IDs that already have attachments imported
    existing_att_ids = _get_existing_attachment_msg_ids()
    log.info("Found %d message IDs with existing attachments", len(existing_att_ids))
    log.info("Already processed in previous runs: %d", len(already_processed))

    # 3. Get token
    token = get_access_token()
    log.info("Token acquired successfully")

    # 4. Get or create the attachment source
    att_source_id = insert_source(
        conn=None,
        source_type="company_email_attachment",
        source_name="email_attachments",
    )

    # 5. Process each email
    processed_count = 0
    total_imported = 0

    for i, doc in enumerate(email_docs):
        if args.limit and processed_count >= args.limit:
            log.info("Limit reached (%d). Stopping.", args.limit)
            break

        msg_id = _extract_graph_msg_id(doc["raw_path"])
        if not msg_id:
            log.debug("Cannot extract msg_id from raw_path: %s", doc["raw_path"])
            continue

        # Skip if already processed in a previous run
        if msg_id in already_processed:
            continue

        processed_count += 1
        stats["checked"] = stats.get("checked", 0) + 1

        log.info(
            "[%d/%d] Checking: %s",
            processed_count,
            min(args.limit, len(email_docs)) if args.limit else len(email_docs),
            doc["title"][:60],
        )

        if args.dry_run:
            # Just list attachments, don't import
            url = f"{GRAPH_BASE}/me/messages/{msg_id}/attachments"
            data = _graph_get_with_retry(
                url, token,
                params={"$select": "id,name,size,contentType"},
                timeout=REQUEST_TIMEOUT_LIST,
                label=f"list-att {doc['title'][:40]}",
            )
            if data and data.get("value"):
                att_names = [a.get("name", "?") for a in data["value"]]
                log.info("  -> %d attachments: %s", len(att_names), ", ".join(att_names))
                stats["with_attachments"] = stats.get("with_attachments", 0) + 1
        else:
            try:
                count = _process_email(doc, msg_id, token, att_source_id)
                if count > 0:
                    stats["with_attachments"] = stats.get("with_attachments", 0) + 1
                    stats["imported"] = stats.get("imported", 0) + count
                    total_imported += count
            except requests.HTTPError as e:
                status = getattr(e.response, "status_code", 0)
                if status in (401, 403):
                    log.error("Auth error — token expired or insufficient permissions. Saving state and exiting.")
                    already_processed.add(msg_id)  # Don't re-add, but save
                    state["processed_msg_ids"] = list(already_processed)
                    state["stats"] = stats
                    _save_state(state)
                    sys.exit(1)
                stats["errors"] = stats.get("errors", 0) + 1
                log.error("HTTP error for %s: %s", doc["title"][:40], e)
            except Exception as e:
                stats["errors"] = stats.get("errors", 0) + 1
                log.error("Unexpected error for %s: %s", doc["title"][:40], e)

        # Mark as processed and save state periodically
        already_processed.add(msg_id)
        if processed_count % 5 == 0:
            state["processed_msg_ids"] = list(already_processed)
            state["stats"] = stats
            _save_state(state)

        # Small delay to avoid throttling
        time.sleep(0.5)

    # Final state save
    state["processed_msg_ids"] = list(already_processed)
    state["stats"] = stats
    state["last_run"] = datetime.now().isoformat()
    _save_state(state)

    # Summary
    log.info("")
    log.info("=" * 60)
    log.info("BATCH COMPLETE")
    log.info("=" * 60)
    log.info("Emails checked this run:    %d", processed_count)
    log.info("Emails with attachments:    %d (cumulative: %d)", total_imported > 0 and 1 or 0, stats.get("with_attachments", 0))
    log.info("Attachments imported:       %d (cumulative: %d)", total_imported, stats.get("imported", 0))
    log.info("Errors:                     %d", stats.get("errors", 0))
    log.info("Total processed (all runs): %d / %d", len(already_processed), len(email_docs))
    log.info("State saved to: %s", STATE_FILE)


if __name__ == "__main__":
    main()
