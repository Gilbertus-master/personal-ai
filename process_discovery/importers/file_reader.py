"""Detect file format and extract text content."""

from __future__ import annotations

import base64
import csv
import json
import subprocess
import xml.etree.ElementTree as ET
from pathlib import Path
from typing import Any

import structlog

log = structlog.get_logger("process_discovery.importers.file_reader")

SUPPORTED_FORMATS: dict[str, str] = {
    ".pdf": "pdf",
    ".docx": "docx", ".doc": "docx",
    ".xlsx": "xlsx", ".xls": "xlsx", ".ods": "xlsx",
    ".pptx": "pptx", ".ppt": "pptx",
    ".csv": "csv",
    ".json": "json", ".jsonl": "json",
    ".png": "image", ".jpg": "image", ".jpeg": "image",
    ".bpmn": "xml", ".xml": "xml",
    ".txt": "text", ".md": "text",
}


def read_file(file_path: str) -> dict[str, Any]:
    """Read any supported file and return extracted text + metadata."""
    path = Path(file_path)
    if not path.exists():
        raise FileNotFoundError(f"File not found: {file_path}")

    ext = path.suffix.lower()
    fmt = SUPPORTED_FORMATS.get(ext)
    if not fmt:
        raise ValueError(
            f"Unsupported format '{ext}'. Supported: {list(SUPPORTED_FORMATS.keys())}"
        )

    readers = {
        "pdf": _read_pdf,
        "docx": _read_docx,
        "xlsx": _read_xlsx,
        "csv": _read_csv,
        "json": _read_json,
        "pptx": _read_pptx,
        "xml": _read_bpmn,
        "image": _read_image,
        "text": _read_text,
    }

    reader = readers[fmt]
    log.info("reading_file", path=str(path), format=fmt)

    raw_text, structured_data, method = reader(str(path))

    return {
        "format": fmt,
        "raw_text": raw_text or "",
        "structured_data": structured_data,
        "pages": _estimate_pages(raw_text),
        "extraction_method": method,
        "filename": path.name,
    }


def _estimate_pages(text: str | None) -> int:
    if not text:
        return 0
    # Rough estimate: ~3000 chars per page
    return max(1, len(text) // 3000)


# ─── Format-specific readers ─────────────────────────────────────────

def _read_pdf(path: str) -> tuple[str, Any, str]:
    """PDF: pdftotext first, fallback to vision OCR."""
    try:
        result = subprocess.run(
            ["pdftotext", path, "-"],
            capture_output=True, text=True, timeout=30,
        )
        if result.stdout.strip():
            return result.stdout.strip(), None, "pdftotext"
    except (FileNotFoundError, subprocess.TimeoutExpired):
        log.warning("pdftotext_unavailable", path=path)

    # Fallback: try pdfminer
    try:
        from pdfminer.high_level import extract_text
        text = extract_text(path)
        if text.strip():
            return text.strip(), None, "pdfminer"
    except ImportError:
        log.warning("pdfminer_unavailable")

    # Last resort: vision OCR (first page only for cost)
    return _read_image(path, media_type="application/pdf")


def _read_docx(path: str) -> tuple[str, Any, str]:
    """DOCX: pandoc first, fallback to python-docx."""
    try:
        result = subprocess.run(
            ["pandoc", path, "-t", "plain", "--wrap=none"],
            capture_output=True, text=True, timeout=30,
        )
        if result.stdout.strip():
            return result.stdout.strip(), None, "pandoc"
    except (FileNotFoundError, subprocess.TimeoutExpired):
        log.warning("pandoc_unavailable")

    # Fallback: python-docx
    try:
        from docx import Document
        doc = Document(path)
        paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
        text = "\n".join(paragraphs)
        return text, None, "python-docx"
    except ImportError:
        raise RuntimeError(
            "Cannot read DOCX: install pandoc (apt install pandoc) "
            "or python-docx (pip install python-docx)"
        )


def _read_xlsx(path: str) -> tuple[str, Any, str]:
    """XLSX: read all sheets as text."""
    try:
        from openpyxl import load_workbook
    except ImportError:
        raise RuntimeError("Install openpyxl: pip install openpyxl")

    wb = load_workbook(path, read_only=True, data_only=True)
    output: list[str] = []
    structured: dict[str, list] = {}

    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        output.append(f"=== Sheet: {sheet_name} ===")
        rows_data: list[list[str]] = []
        for row in ws.iter_rows(values_only=True):
            if any(cell is not None for cell in row):
                row_strs = [str(c) if c is not None else "" for c in row]
                output.append(" | ".join(row_strs))
                rows_data.append(row_strs)
        structured[sheet_name] = rows_data

    wb.close()
    return "\n".join(output), structured, "openpyxl"


def _read_csv(path: str) -> tuple[str, Any, str]:
    """CSV: read as text table."""
    with open(path, newline="", encoding="utf-8-sig") as f:
        reader = csv.reader(f)
        rows = list(reader)

    if not rows:
        return "", None, "csv"

    output = []
    for row in rows:
        output.append(" | ".join(row))
    return "\n".join(output), rows, "csv"


def _read_json(path: str) -> tuple[str, Any, str]:
    """JSON: pretty-print."""
    with open(path, encoding="utf-8") as f:
        data = json.load(f)
    text = json.dumps(data, ensure_ascii=False, indent=2)
    return text, data, "json"


def _read_pptx(path: str) -> tuple[str, Any, str]:
    """PPTX: extract text from slides."""
    try:
        from pptx import Presentation
    except ImportError:
        raise RuntimeError("Install python-pptx: pip install python-pptx")

    prs = Presentation(path)
    output: list[str] = []

    for i, slide in enumerate(prs.slides, 1):
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame and shape.text.strip():
                texts.append(shape.text.strip())
        if texts:
            output.append(f"=== Slide {i} ===\n" + "\n".join(texts))

    return "\n\n".join(output), None, "python-pptx"


def _read_bpmn(path: str) -> tuple[str, Any, str]:
    """BPMN/XML: extract process elements and flows."""
    tree = ET.parse(path)
    root = tree.getroot()

    # Strip namespace for easier querying
    ns = ""
    if root.tag.startswith("{"):
        ns = root.tag.split("}")[0] + "}"

    output: list[str] = []
    elements: list[dict] = []

    # Find process elements
    for proc in root.iter(f"{ns}process"):
        proc_name = proc.get("name", proc.get("id", "unnamed"))
        output.append(f"=== Process: {proc_name} ===")

        for tag in ("task", "userTask", "serviceTask", "manualTask",
                     "sendTask", "receiveTask", "scriptTask",
                     "startEvent", "endEvent", "exclusiveGateway",
                     "parallelGateway", "intermediateThrowEvent"):
            for elem in proc.iter(f"{ns}{tag}"):
                name = elem.get("name", elem.get("id", ""))
                if name:
                    output.append(f"  [{tag}] {name}")
                    elements.append({"type": tag, "name": name})

        # Sequence flows
        for flow in proc.iter(f"{ns}sequenceFlow"):
            src = flow.get("sourceRef", "?")
            tgt = flow.get("targetRef", "?")
            output.append(f"  {src} → {tgt}")

    if not output:
        # Generic XML fallback
        output.append(ET.tostring(root, encoding="unicode", method="text")[:10000])

    return "\n".join(output), elements, "bpmn_xml"


def _read_image(path: str, media_type: str | None = None) -> tuple[str, Any, str]:
    """Image/PDF: OCR via Claude vision API."""
    from anthropic import Anthropic
    from app.db.cost_tracker import log_anthropic_cost

    with open(path, "rb") as f:
        image_data = base64.b64encode(f.read()).decode()

    if not media_type:
        ext = Path(path).suffix.lower()
        media_type = {
            ".jpg": "image/jpeg", ".jpeg": "image/jpeg",
            ".png": "image/png", ".pdf": "application/pdf",
        }.get(ext, "image/png")

    client = Anthropic()
    model = "claude-haiku-4-5-20251001"

    response = client.messages.create(
        model=model,
        max_tokens=4096,
        messages=[{
            "role": "user",
            "content": [
                {
                    "type": "image",
                    "source": {
                        "type": "base64",
                        "media_type": media_type,
                        "data": image_data,
                    },
                },
                {
                    "type": "text",
                    "text": (
                        "Wyciągnij z tego diagramu/dokumentu wszystkie procesy biznesowe. "
                        "Dla każdego procesu podaj: nazwę, opis, kroki/etapy, "
                        "kto jest odpowiedzialny (jeśli widoczne). "
                        "Odpowiedz w formacie tekstowym, każdy proces w oddzielnej sekcji."
                    ),
                },
            ],
        }],
    )

    if hasattr(response, "usage"):
        log_anthropic_cost(model, "process_import_ocr", response.usage)

    return response.content[0].text, None, "claude_vision"


def _read_text(path: str) -> tuple[str, Any, str]:
    """Plain text / Markdown."""
    with open(path, encoding="utf-8") as f:
        text = f.read()
    return text, None, "text"
